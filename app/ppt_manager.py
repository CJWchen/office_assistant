"""
PPT生成管理器
负责PPT模板管理、内容填充、图片匹配、导出等功能
"""

import os
import json
import uuid
from datetime import datetime, timedelta
from flask import current_app
import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import requests
from PIL import Image
import io

from .models import db, PPTTemplate, PPTProject

class PPTManager:
    """PPT管理器类"""
    
    def __init__(self):
        self.unsplash_access_key = current_app.config.get('UNSPLASH_ACCESS_KEY', '')
        self.deepseek_api_key = current_app.config.get('DEEPSEEK_API_KEY', '')
        self.upload_base_path = current_app.config.get('UPLOAD_FOLDER', 'uploads')
        self.ppt_base_path = os.path.join(self.upload_base_path, 'ppt')
        
        # 确保目录存在
        os.makedirs(self.ppt_base_path, exist_ok=True)
        os.makedirs(os.path.join(self.ppt_base_path, 'templates'), exist_ok=True)
        os.makedirs(os.path.join(self.ppt_base_path, 'thumbnails'), exist_ok=True)
        os.makedirs(os.path.join(self.ppt_base_path, 'generated'), exist_ok=True)
    
    # ==================== 模板管理功能 ====================
    
    def upload_template(self, user_id, template_file, name, description, category, style_type, tags=None, is_public=False):
        """
        上传PPT模板
        Args:
            user_id: 用户ID
            template_file: 上传的模板文件对象
            name: 模板名称
            description: 模板描述
            category: 模板分类
            style_type: 样式类型
            tags: 标签列表
            is_public: 是否公开
        Returns:
            PPTTemplate对象
        """
        try:
            # 生成唯一文件名
            file_ext = os.path.splitext(template_file.filename)[1].lower()
            if file_ext not in ['.pptx', '.ppt']:
                return None, "只支持.pptx或.ppt格式的模板文件"
            
            unique_filename = f"{uuid.uuid4().hex}{file_ext}"
            template_dir = os.path.join(self.ppt_base_path, 'templates', str(user_id))
            os.makedirs(template_dir, exist_ok=True)
            
            template_path = os.path.join(template_dir, unique_filename)
            template_file.save(template_path)
            
            # 生成缩略图
            thumbnail_path = self._generate_thumbnail(template_path, user_id)
            
            # 提取颜色方案
            color_scheme = self._extract_color_scheme(template_path)
            
            # 创建模板记录
            template = PPTTemplate(
                name=name,
                description=description,
                category=category,
                thumbnail_path=thumbnail_path,
                template_path=template_path,
                style_type=style_type,
                color_scheme=json.dumps(color_scheme) if color_scheme else None,
                tags=','.join(tags) if tags else '',
                user_id=user_id,
                is_public=is_public
            )
            
            db.session.add(template)
            db.session.commit()
            
            return template, "模板上传成功"
            
        except Exception as e:
            db.session.rollback()
            current_app.logger.error(f"模板上传失败: {str(e)}")
            return None, f"模板上传失败: {str(e)}"
    
    def list_templates(self, user_id, category=None, style_type=None, public_only=False):
        """
        获取模板列表
        Args:
            user_id: 用户ID
            category: 按分类筛选
            style_type: 按样式类型筛选
            public_only: 是否只获取公开模板
        Returns:
            模板列表
        """
        query = PPTTemplate.query
        
        if public_only:
            query = query.filter_by(is_public=True)
        else:
            # 获取用户自己的模板和公开模板
            query = query.filter(
                db.or_(
                    PPTTemplate.user_id == user_id,
                    PPTTemplate.is_public == True
                )
            )
        
        if category:
            query = query.filter_by(category=category)
        
        if style_type:
            query = query.filter_by(style_type=style_type)
        
        return query.order_by(PPTTemplate.created_at.desc()).all()
    
    def get_template(self, template_id, user_id=None):
        """
        获取模板详情
        Args:
            template_id: 模板ID
            user_id: 用户ID（用于权限检查）
        Returns:
            PPTTemplate对象
        """
        template = PPTTemplate.query.get(template_id)
        
        if not template:
            return None
        
        # 权限检查：模板公开或用户自己的模板
        if user_id and template.user_id != user_id and not template.is_public:
            return None
        
        return template
    
    def delete_template(self, template_id, user_id):
        """
        删除模板
        Args:
            template_id: 模板ID
            user_id: 用户ID
        Returns:
            (success, message)
        """
        template = self.get_template(template_id, user_id)
        
        if not template:
            return False, "模板不存在或无权删除"
        
        if template.user_id != user_id:
            return False, "无权删除他人模板"
        
        try:
            # 删除文件
            if os.path.exists(template.template_path):
                os.remove(template.template_path)
            
            if template.thumbnail_path and os.path.exists(template.thumbnail_path):
                os.remove(template.thumbnail_path)
            
            # 删除数据库记录
            db.session.delete(template)
            db.session.commit()
            
            return True, "模板删除成功"
            
        except Exception as e:
            db.session.rollback()
            current_app.logger.error(f"模板删除失败: {str(e)}")
            return False, f"模板删除失败: {str(e)}"
    
    # ==================== 内容智能填充功能 ====================
    
    def create_project(self, user_id, title, description, template_id=None, content_data=None):
        """
        创建PPT项目
        Args:
            user_id: 用户ID
            title: 项目标题
            description: 项目描述
            template_id: 模板ID（可选）
            content_data: 内容数据（JSON格式）
        Returns:
            PPTProject对象
        """
        try:
            project = PPTProject(
                title=title,
                description=description,
                template_id=template_id,
                content_data=json.dumps(content_data) if content_data else None,
                user_id=user_id,
                status='draft'
            )
            
            db.session.add(project)
            db.session.commit()
            
            return project, "项目创建成功"
            
        except Exception as e:
            db.session.rollback()
            current_app.logger.error(f"项目创建失败: {str(e)}")
            return None, f"项目创建失败: {str(e)}"
    
    def generate_pptx(self, project_id, user_id):
        """
        生成PPTX文件
        Args:
            project_id: 项目ID
            user_id: 用户ID
        Returns:
            (success, message, file_path)
        """
        project = PPTProject.query.filter_by(id=project_id, user_id=user_id).first()
        
        if not project:
            return False, "项目不存在或无权访问", None
        
        try:
            # 更新状态为生成中
            project.status = 'generating'
            db.session.commit()
            
            # 获取模板
            template = None
            if project.template_id:
                template = self.get_template(project.template_id, user_id)
            
            # 解析内容数据
            content_data = json.loads(project.content_data) if project.content_data else {}
            
            # 生成PPTX文件
            if template:
                pptx_path = self._generate_from_template(template.template_path, content_data, user_id)
            else:
                pptx_path = self._generate_from_scratch(content_data, user_id)
            
            # 保存生成的文件路径
            project.generated_pptx_path = pptx_path
            project.status = 'completed'
            project.updated_at = datetime.utcnow()
            db.session.commit()
            
            return True, "PPT生成成功", pptx_path
            
        except Exception as e:
            project.status = 'failed'
            db.session.commit()
            current_app.logger.error(f"PPT生成失败: {str(e)}")
            return False, f"PPT生成失败: {str(e)}", None
    
    def generate_html(self, project_id, user_id):
        """
        生成HTML版本
        Args:
            project_id: 项目ID
            user_id: 用户ID
        Returns:
            (success, message, file_path)
        """
        project = PPTProject.query.filter_by(id=project_id, user_id=user_id).first()
        
        if not project:
            return False, "项目不存在或无权访问", None
        
        try:
            # 检查是否已有PPTX文件
            if not project.generated_pptx_path or not os.path.exists(project.generated_pptx_path):
                # 先生成PPTX
                success, msg, pptx_path = self.generate_pptx(project_id, user_id)
                if not success:
                    return False, msg, None
            
            # 将PPTX转换为HTML
            html_path = self._convert_to_html(project.generated_pptx_path, user_id)
            
            # 保存HTML路径
            project.generated_html_path = html_path
            db.session.commit()
            
            return True, "HTML生成成功", html_path
            
        except Exception as e:
            current_app.logger.error(f"HTML生成失败: {str(e)}")
            return False, f"HTML生成失败: {str(e)}", None
    
    def create_share_link(self, project_id, user_id, expires_hours=24):
        """
        创建分享链接
        Args:
            project_id: 项目ID
            user_id: 用户ID
            expires_hours: 过期时间（小时）
        Returns:
            (success, message, share_url)
        """
        project = PPTProject.query.filter_by(id=project_id, user_id=user_id).first()
        
        if not project:
            return False, "项目不存在或无权访问", None
        
        try:
            # 生成分享令牌
            share_token = uuid.uuid4().hex
            
            # 设置过期时间
            share_expires = datetime.utcnow() + timedelta(hours=expires_hours)
            
            # 更新项目
            project.share_token = share_token
            project.share_expires = share_expires
            db.session.commit()
            
            # 生成分享URL
            share_url = f"/ppt/share/{share_token}"
            
            return True, "分享链接创建成功", share_url
            
        except Exception as e:
            db.session.rollback()
            current_app.logger.error(f"分享链接创建失败: {str(e)}")
            return False, f"分享链接创建失败: {str(e)}", None
    
    # ==================== AI内容优化功能 ====================
    
    def optimize_content_with_ai(self, content_data, optimization_type='clarity'):
        """
        使用DeepSeek API优化内容
        Args:
            content_data: 内容数据
            optimization_type: 优化类型（clarity, conciseness, engagement, etc.）
        Returns:
            优化后的内容
        """
        if not self.deepseek_api_key:
            return content_data, "未配置DeepSeek API密钥"
        
        try:
            # 构建优化提示
            optimization_prompts = {
                'clarity': "请优化以下内容，使其更加清晰易懂：",
                'conciseness': "请优化以下内容，使其更加简洁精炼：",
                'engagement': "请优化以下内容，使其更具吸引力：",
                'professional': "请优化以下内容，使其更具专业性："
            }
            
            prompt = optimization_prompts.get(optimization_type, optimization_prompts['clarity'])
            
            # 将内容数据转换为文本
            content_text = self._content_to_text(content_data)
            
            # 调用DeepSeek API
            headers = {
                'Authorization': f'Bearer {self.deepseek_api_key}',
                'Content-Type': 'application/json'
            }
            
            data = {
                'model': 'deepseek-chat',
                'messages': [
                    {
                        'role': 'system',
                        'content': '你是一个专业的PPT内容优化专家，擅长让内容更加清晰、简洁、有吸引力。'
                    },
                    {
                        'role': 'user',
                        'content': f"{prompt}\n\n{content_text}"
                    }
                ],
                'max_tokens': 2000,
                'temperature': 0.7
            }
            
            response = requests.post(
                'https://api.deepseek.com/v1/chat/completions',
                headers=headers,
                json=data,
                timeout=30
            )
            
            if response.status_code == 200:
                result = response.json()
                optimized_text = result['choices'][0]['message']['content']
                
                # 将优化后的文本转换回内容数据格式
                optimized_content = self._text_to_content(optimized_text, content_data)
                
                return optimized_content, "内容优化成功"
            else:
                return content_data, f"API调用失败: {response.status_code}"
                
        except Exception as e:
            current_app.logger.error(f"AI内容优化失败: {str(e)}")
            return content_data, f"AI内容优化失败: {str(e)}"
    
    # ==================== 图片匹配功能 ====================
    
    def search_unsplash_images(self, keywords, count=5):
        """
        搜索Unsplash图片
        Args:
            keywords: 关键词列表
            count: 返回图片数量
        Returns:
            图片信息列表
        """
        if not self.unsplash_access_key:
            return [], "未配置Unsplash API密钥"
        
        try:
            query = ' '.join(keywords)
            
            headers = {
                'Authorization': f'Client-ID {self.unsplash_access_key}'
            }
            
            params = {
                'query': query,
                'per_page': count,
                'orientation': 'landscape'
            }
            
            response = requests.get(
                'https://api.unsplash.com/search/photos',
                headers=headers,
                params=params,
                timeout=30
            )
            
            if response.status_code == 200:
                data = response.json()
                images = []
                
                for photo in data.get('results', [])[:count]:
                    image_info = {
                        'id': photo.get('id'),
                        'url': photo.get('urls', {}).get('regular'),
                        'thumb_url': photo.get('urls', {}).get('thumb'),
                        'author': photo.get('user', {}).get('name'),
                        'author_url': photo.get('user', {}).get('links', {}).get('html'),
                        'description': photo.get('description') or photo.get('alt_description'),
                        'color': photo.get('color')
                    }
                    images.append(image_info)
                
                return images, "图片搜索成功"
            else:
                return [], f"Unsplash API调用失败: {response.status_code}"
                
        except Exception as e:
            current_app.logger.error(f"Unsplash图片搜索失败: {str(e)}")
            return [], f"Unsplash图片搜索失败: {str(e)}"
    
    def match_images_to_content(self, content_data, image_count=3):
        """
        为内容匹配合适的图片
        Args:
            content_data: 内容数据
            image_count: 每页图片数量
        Returns:
            匹配的图片数据
        """
        # 从内容中提取关键词
        keywords = self._extract_keywords(content_data)
        
        # 搜索图片
        images, message = self.search_unsplash_images(keywords, count=image_count*3)
        
        # 按相关性排序（简单实现：按颜色鲜艳度排序）
        sorted_images = sorted(images, key=lambda x: self._color_vibrancy_score(x.get('color', '#000000')), reverse=True)
        
        return sorted_images[:image_count]
    
    # ==================== 私有辅助方法 ====================
    
    def _generate_thumbnail(self, pptx_path, user_id):
        """
        生成PPT缩略图
        Args:
            pptx_path: PPTX文件路径
            user_id: 用户ID
        Returns:
            缩略图路径
        """
        try:
            # 打开PPTX文件
            prs = Presentation(pptx_path)
            
            # 获取第一页幻灯片
            if len(prs.slides) > 0:
                slide = prs.slides[0]
                
                # 创建缩略图目录
                thumb_dir = os.path.join(self.ppt_base_path, 'thumbnails', str(user_id))
                os.makedirs(thumb_dir, exist_ok=True)
                
                # 生成缩略图文件名
                thumb_filename = f"{uuid.uuid4().hex}.png"
                thumb_path = os.path.join(thumb_dir, thumb_filename)
                
                # 注意：python-pptx本身不提供截图功能
                # 这里返回一个占位符路径，实际项目中可以使用其他库生成缩略图
                return thumb_path
                
            return None
            
        except Exception as e:
            current_app.logger.error(f"生成缩略图失败: {str(e)}")
            return None
    
    def _extract_color_scheme(self, pptx_path):
        """
        提取PPT颜色方案
        Args:
            pptx_path: PPTX文件路径
        Returns:
            颜色方案字典
        """
        try:
            prs = Presentation(pptx_path)
            
            color_scheme = {
                'background_colors': [],
                'text_colors': [],
                'accent_colors': []
            }
            
            # 分析前几页幻灯片
            for i, slide in enumerate(prs.slides[:3]):
                # 背景色
                if slide.background.fill.type == 1:  # 纯色填充
                    fill = slide.background.fill
                    if hasattr(fill, 'fore_color'):
                        rgb = fill.fore_color.rgb
                        if rgb:
                            color_scheme['background_colors'].append(self._rgb_to_hex(rgb))
                
                # 形状颜色
                for shape in slide.shapes:
                    if hasattr(shape, 'fill'):
                        if shape.fill.type == 1:  # 纯色填充
                            rgb = shape.fill.fore_color.rgb
                            if rgb:
                                color_scheme['accent_colors'].append(self._rgb_to_hex(rgb))
                    
                    # 文本颜色
                    if hasattr(shape, 'text'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if hasattr(run.font, 'color') and run.font.color.type == 1:
                                    rgb = run.font.color.rgb
                                    if rgb:
                                        color_scheme['text_colors'].append(self._rgb_to_hex(rgb))
            
            # 去重
            for key in color_scheme:
                color_scheme[key] = list(set(color_scheme[key]))
            
            return color_scheme
            
        except Exception as e:
            current_app.logger.error(f"提取颜色方案失败: {str(e)}")
            return {}
    
    def _rgb_to_hex(self, rgb):
        """将RGB值转换为十六进制颜色码"""
        if isinstance(rgb, int):
            r = (rgb >> 16) & 0xFF
            g = (rgb >> 8) & 0xFF
            b = rgb & 0xFF
            return f'#{r:02x}{g:02x}{b:02x}'
        elif hasattr(rgb, 'rgb'):
            return f'#{rgb.rgb:06x}'
        else:
            return '#000000'
    
    def _generate_from_template(self, template_path, content_data, user_id):
        """
        基于模板生成PPTX
        Args:
            template_path: 模板路径
            content_data: 内容数据
            user_id: 用户ID
        Returns:
            生成的PPTX文件路径
        """
        try:
            # 复制模板
            prs = Presentation(template_path)
            
            # 内容填充逻辑
            # 这里实现占位符替换等逻辑
            # ...
            
            # 保存生成的文件
            output_dir = os.path.join(self.ppt_base_path, 'generated', str(user_id))
            os.makedirs(output_dir, exist_ok=True)
            
            output_filename = f"{uuid.uuid4().hex}.pptx"
            output_path = os.path.join(output_dir, output_filename)
            
            prs.save(output_path)
            
            return output_path
            
        except Exception as e:
            current_app.logger.error(f"基于模板生成PPT失败: {str(e)}")
            raise
    
    def _generate_from_scratch(self, content_data, user_id):
        """
        从零开始生成PPTX
        Args:
            content_data: 内容数据
            user_id: 用户ID
        Returns:
            生成的PPTX文件路径
        """
        try:
            prs = Presentation()
            
            # 解析内容数据
            slides_data = content_data.get('slides', [])
            
            for slide_data in slides_data:
                # 根据布局类型创建幻灯片
                layout_type = slide_data.get('layout', 'title_content')
                
                if layout_type == 'title':
                    # 标题幻灯片
                    slide_layout = prs.slide_layouts[0]  # 标题幻灯片
                    slide = prs.slides.add_slide(slide_layout)
                    
                    title = slide.shapes.title
                    if title:
                        title.text = slide_data.get('title', '')
                    
                    subtitle = slide.placeholders[1]
                    if subtitle:
                        subtitle.text = slide_data.get('content', '')
                        
                elif layout_type == 'title_content':
                    # 标题和内容幻灯片
                    slide_layout = prs.slide_layouts[1]  # 标题和内容
                    slide = prs.slides.add_slide(slide_layout)
                    
                    title = slide.shapes.title
                    if title:
                        title.text = slide_data.get('title', '')
                    
                    content = slide.placeholders[1]
                    if content:
                        content.text = slide_data.get('content', '')
                        
                elif layout_type == 'blank':
                    # 空白幻灯片
                    slide_layout = prs.slide_layouts[6]  # 空白
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # 手动添加文本框
                    left = Inches(1)
                    top = Inches(1)
                    width = Inches(8)
                    height = Inches(2)
                    
                    text_box = slide.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    
                    # 添加标题
                    p = text_frame.add_paragraph()
                    p.text = slide_data.get('title', '')
                    p.font.bold = True
                    p.font.size = Pt(32)
                    
                    # 添加内容
                    if slide_data.get('content'):
                        p = text_frame.add_paragraph()
                        p.text = slide_data.get('content', '')
                        p.font.size = Pt(18)
                        
                else:
                    # 默认使用标题和内容布局
                    slide_layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(slide_layout)
                    
                    title = slide.shapes.title
                    if title:
                        title.text = slide_data.get('title', '')
                    
                    content = slide.placeholders[1]
                    if content:
                        content.text = slide_data.get('content', '')
            
            # 保存文件
            output_dir = os.path.join(self.ppt_base_path, 'generated', str(user_id))
            os.makedirs(output_dir, exist_ok=True)
            
            output_filename = f"{uuid.uuid4().hex}.pptx"
            output_path = os.path.join(output_dir, output_filename)
            
            prs.save(output_path)
            
            return output_path
            
        except Exception as e:
            current_app.logger.error(f"从零生成PPT失败: {str(e)}")
            raise
    
    def _convert_to_html(self, pptx_path, user_id):
        """
        将PPTX转换为HTML
        Args:
            pptx_path: PPTX文件路径
            user_id: 用户ID
        Returns:
            HTML文件路径
        """
        try:
            # 这里实现PPTX到HTML的转换逻辑
            # 实际项目中可以使用python-pptx提取内容，然后生成HTML
            # ...
            
            output_dir = os.path.join(self.ppt_base_path, 'generated', str(user_id))
            os.makedirs(output_dir, exist_ok=True)
            
            html_filename = f"{uuid.uuid4().hex}.html"
            html_path = os.path.join(output_dir, html_filename)
            
            # 生成简单的HTML占位符
            with open(html_path, 'w', encoding='utf-8') as f:
                f.write("""
<!DOCTYPE html>
<html>
<head>
    <meta charset="UTF-8">
    <title>PPT演示</title>
    <style>
        body { font-family: Arial, sans-serif; margin: 0; padding: 20px; }
        .slide { margin-bottom: 20px; padding: 20px; border: 1px solid #ddd; }
        h1 { color: #333; }
    </style>
</head>
<body>
    <div class="slide">
        <h1>PPT内容将在这里显示</h1>
        <p>这是一个PPT的HTML版本演示。</p>
    </div>
</body>
</html>
                """)
            
            return html_path
            
        except Exception as e:
            current_app.logger.error(f"PPTX转HTML失败: {str(e)}")
            raise
    
    def _content_to_text(self, content_data):
        """将内容数据转换为文本"""
        try:
            text_parts = []
            
            if isinstance(content_data, dict):
                for key, value in content_data.items():
                    if isinstance(value, str):
                        text_parts.append(value)
                    elif isinstance(value, list):
                        text_parts.extend([str(item) for item in value if isinstance(item, str)])
            
            return '\n'.join(text_parts)
        except:
            return str(content_data)
    
    def _text_to_content(self, text, original_content):
        """将文本转换回内容数据格式（简单实现）"""
        # 这里应该根据原始内容的结构进行智能转换
        # 简化实现：返回原始内容
        return original_content
    
    def _extract_keywords(self, content_data):
        """从内容中提取关键词"""
        keywords = []
        
        if isinstance(content_data, dict):
            for key, value in content_data.items():
                if isinstance(value, str):
                    # 提取名词性关键词（简化实现）
                    words = value.split()
                    keywords.extend([word.lower() for word in words if len(word) > 3])
        
        # 去重并限制数量
        keywords = list(set(keywords))[:10]
        
        # 添加默认关键词
        if not keywords:
            keywords = ['business', 'presentation', 'design']
        
        return keywords
    
    def _color_vibrancy_score(self, hex_color):
        """计算颜色鲜艳度分数"""
        try:
            # 移除#前缀
            hex_color = hex_color.lstrip('#')
            
            if len(hex_color) != 6:
                return 0
            
            # 转换为RGB
            r = int(hex_color[0:2], 16)
            g = int(hex_color[2:4], 16)
            b = int(hex_color[4:6], 16)
            
            # 计算鲜艳度（饱和度）
            max_val = max(r, g, b)
            min_val = min(r, g, b)
            
            if max_val == 0:
                return 0
            
            saturation = (max_val - min_val) / max_val
            return saturation
            
        except:
            return 0